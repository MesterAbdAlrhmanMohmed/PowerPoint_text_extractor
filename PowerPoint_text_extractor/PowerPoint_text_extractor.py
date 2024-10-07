from PyQt6 import QtWidgets as qt
from PyQt6 import QtGui as qt1
from PyQt6 import QtCore as qt2
from PyQt6.QtPrintSupport import QPrinter, QPrintDialog
from PyQt6.QtCore import Qt
from pptx import Presentation
import about,winsound,pyperclip,user_guide
class PowerPointTextExtractor(qt.QMainWindow):
    def __init__(self):
        super().__init__()
        qt1.QShortcut("ctrl+=", self).activated.connect(self.increase_font_size)
        qt1.QShortcut("ctrl+-", self).activated.connect(self.decrease_font_size)
        qt1.QShortcut("ctrl+c", self).activated.connect(self.copy_line)
        qt1.QShortcut("ctrl+a", self).activated.connect(self.copy_text)
        qt1.QShortcut("ctrl+p", self).activated.connect(self.print_text)
        qt1.QShortcut("ctrl+s", self).activated.connect(self.save_text_as_txt)
        self.setWindowTitle("PowerPoint Text Extractor")
        self.resize(1000,500)
        self.choose_file_btn=qt.QPushButton("إختيار ملف PowerPoint")
        self.choose_file_btn.setDefault(True)
        self.choose_file_btn.clicked.connect(self.choose_PowerPoint_file)
        self.show_path_label=qt.QLabel("مسار ملف PowerPoint المحدد")
        self.file_path=qt.QLineEdit()
        self.file_path.setAccessibleName("مسار ملف PowerPoint المحدد")
        self.file_path.setReadOnly(True)
        self.start_extraction_btn=qt.QPushButton("بدء استخراج النص")
        self.start_extraction_btn.setDefault(True)
        self.start_extraction_btn.clicked.connect(self.extract_text)
        self.text_edit=qt.QTextEdit()
        self.text_edit.setReadOnly(True)
        self.text_edit.setTextInteractionFlags(
        Qt.TextInteractionFlag.TextSelectableByKeyboard | Qt.TextInteractionFlag.TextSelectableByMouse)
        self.text_edit.setLineWrapMode(qt.QTextEdit.LineWrapMode.NoWrap)
        self.text_edit.setAccessibleName("النص المستخرج")
        self.about_btn=qt.QPushButton("عن المطور")
        self.about_btn.setDefault(True)
        self.about_btn.clicked.connect(self.about)        
        self.UserGuide=qt.QPushButton("دليل المستخدم")
        self.UserGuide.setDefault(True)
        self.UserGuide.clicked.connect(self.guide)
        self.font_size=20
        font=self.font()
        font.setPointSize(self.font_size)
        self.text_edit.setFont(font)        
        layout=qt.QVBoxLayout()        
        layout.addWidget(self.choose_file_btn)
        layout.addWidget(self.show_path_label)
        layout.addWidget(self.file_path)
        layout.addWidget(self.start_extraction_btn)
        layout.addWidget(self.text_edit)
        layout.addWidget(self.about_btn)        
        layout.addWidget(self.UserGuide)
        container=qt.QWidget()
        container.setLayout(layout)
        self.setCentralWidget(container)
    def guide(self):
        user_guide.dialog(self).exec()
    def choose_PowerPoint_file(self):
        file_dialog=qt.QFileDialog()
        file_name, _ = file_dialog.getOpenFileName(self, "اختر ملف PowerPoint", "", "PowerPoint Files (*.pptx)")
        if file_name:
            self.file_path.setText(file_name)
    def extract_text(self):
        try:
            ppt=Presentation(self.file_path.text())
            extracted_text=""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
            self.text_edit.setText(extracted_text)
            self.text_edit.setFocus()
        except Exception as error:
            qt.QMessageBox.critical(self, "خطأ في استخراج النص", str(error))
    def about(self):
        about.dialog(self).exec()
    def increase_font_size(self):
        self.font_size += 1
        self.update_font_size()
    def decrease_font_size(self):
        self.font_size -= 1
        self.update_font_size()
    def update_font_size(self):
        cursor=self.text_edit.textCursor()
        self.text_edit.selectAll()
        font=self.text_edit.font()
        font.setPointSize(self.font_size)
        self.text_edit.setCurrentFont(font)
        self.text_edit.setTextCursor(cursor)
    def copy_line(self):
        try:
            cursor=self.text_edit.textCursor()
            if cursor.hasSelection():
                selected_text = cursor.selectedText()
                pyperclip.copy(selected_text)
                winsound.Beep(1000, 100)
        except Exception as error:
            qt.QMessageBox.critical(self, "تنبيه حدث خطأ", str(error))
    def copy_text(self):
        try:
            text_content=self.text_edit.toPlainText()
            pyperclip.copy(text_content)
            winsound.Beep(1000, 100)
        except Exception as error:
            qt.QMessageBox.critical(self, "تنبيه حدث خطأ", str(error))
    def print_text(self):
        try:
            printer=QPrinter()
            dialog=QPrintDialog(printer, self)
            if dialog.exec() == QPrintDialog.DialogCode.Accepted:
                self.text_edit.print(printer)
        except Exception as error:
            qt.QMessageBox.critical(self, "تنبيه حدث خطأ", str(error))
    def save_text_as_txt(self):
        try:
            file_dialog=qt.QFileDialog()
            file_dialog.setAcceptMode(qt.QFileDialog.AcceptMode.AcceptSave)
            file_dialog.setNameFilter("Text Files (*.txt);;All Files (*)")
            file_dialog.setDefaultSuffix("txt")
            if file_dialog.exec() == qt.QFileDialog.DialogCode.Accepted:
                file_name=file_dialog.selectedFiles()[0]
                with open(file_name, 'w', encoding='utf-8') as file:
                    text_content=self.text_edit.toPlainText()
                    file.write(text_content)
        except Exception as error:
            qt.QMessageBox.critical(self, "تنبيه حدث خطأ", str(error))
app=qt.QApplication([])            
app.setStyle("fusion")
window=PowerPointTextExtractor()
window.show()
app.exec()